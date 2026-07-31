import io
from typing import List, Optional
import pptx

from backend.app.services.content_extraction.extractors.base_extractor import (
    BaseContentExtractor, ExtractedBlockData
)
from backend.app.services.content_extraction.prompt_injection_detector import PromptInjectionDetector

class PPTXContentExtractor(BaseContentExtractor):
    async def extract_blocks(
        self,
        file_bytes: bytes,
        document_id: str,
        document_version_id: str,
        lesson_id: Optional[str] = None
    ) -> List[ExtractedBlockData]:
        extracted_blocks: List[ExtractedBlockData] = []
        prs = pptx.Presentation(io.BytesIO(file_bytes))
        seq = 1

        for slide_idx, slide in enumerate(prs.slides):
            slide_num = slide_idx + 1
            try:
                slide_text_parts = []
                title_text = ""

                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text_frame_str = shape.text_frame.text.strip()
                        if text_frame_str:
                            if shape == slide.shapes.title and not title_text:
                                title_text = text_frame_str
                            else:
                                slide_text_parts.append(text_frame_str)

                # Emit title heading if present
                if title_text:
                    has_inj, inj_signals = PromptInjectionDetector.detect(title_text)
                    extracted_blocks.append(
                        ExtractedBlockData(
                            block_type="heading",
                            raw_content=title_text,
                            normalized_content=title_text,
                            language="vi",
                            slide_number=slide_num,
                            sequence_number=seq,
                            metadata={"has_prompt_injection": has_inj, "prompt_injection_signals": inj_signals},
                            extraction_confidence=1.0,
                            source_reference=f"Slide {slide_num}"
                        )
                    )
                    seq += 1

                # Emit slide content body
                if slide_text_parts:
                    body_text = "\n".join(slide_text_parts)
                    has_inj, inj_signals = PromptInjectionDetector.detect(body_text)
                    
                    block_type = "paragraph"
                    if any(line.strip().startswith(("-", "*", "1.", "2.")) for line in slide_text_parts):
                        block_type = "list"

                    extracted_blocks.append(
                        ExtractedBlockData(
                            block_type=block_type,
                            raw_content=body_text,
                            normalized_content=body_text,
                            language="vi",
                            slide_number=slide_num,
                            sequence_number=seq,
                            metadata={"has_prompt_injection": has_inj, "prompt_injection_signals": inj_signals},
                            extraction_confidence=0.95,
                            source_reference=f"Slide {slide_num}"
                        )
                    )
                    seq += 1
            except Exception as slide_err:
                # Isolated slide failure handling (Req 12)
                extracted_blocks.append(
                    ExtractedBlockData(
                        block_type="note",
                        raw_content="",
                        normalized_content=f"[ERR_SLIDE_EXTRACTION_FAILED: Slide {slide_num}] {str(slide_err)}",
                        language="vi",
                        slide_number=slide_num,
                        sequence_number=seq,
                        metadata={"error": str(slide_err), "partial_failure": True},
                        extraction_confidence=0.1,
                        source_reference=f"Slide {slide_num}"
                    )
                )
                seq += 1

        return extracted_blocks
